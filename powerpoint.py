from pptx import Presentation
from pptx.util import Inches

# Création d'une nouvelle présentation PowerPoint
prs = Presentation()

# Ajout d'une diapositive de titre
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Présentation de mon script Python"
subtitle.text = "Démonstration de son fonctionnement"

# Ajout du contenu de votre script à la présentation
code_snippet = """
import openai
import discord
import yaml
import json
import requests
from components.Brave import brave_api, extract_descriptions_and_urls_to_json
from components.agents import chatgpt_reply
# Votre script continue ici...
"""
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Code Python"
content.text = code_snippet

# Enregistrement de la présentation
prs.save("presentation.pptx")
